from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, font_size, font_color=RGBColor(255, 255, 255), bold=False, font_name="Segoe UI", alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, hollow=False):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    if hollow:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.color.rgb = fill_color
    return shape

BG_COLOR = RGBColor(15, 18, 25) # Dark Blue/Black
CARD_BG = RGBColor(28, 33, 45) # Lighter Dark Blue
CYAN = RGBColor(0, 229, 255) # Neon Cyan
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(170, 175, 185)

prs = Presentation()
blank_layout = prs.slide_layouts[6]

# ---------------------------------------------------------
# SLIDE 1: Capa
# ---------------------------------------------------------
slide1 = prs.slides.add_slide(blank_layout)
set_background(slide1, BG_COLOR)
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.2), CYAN)
add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.6), Inches(0.1), Inches(2.5), CYAN)

add_text(slide1, "SISTEMA DE", Inches(0.8), Inches(2.5), Inches(8), Inches(1), Pt(32), GRAY, bold=True)
add_text(slide1, "QUIMIOINFORMÁTICA\nINTELIGENTE", Inches(0.8), Inches(3), Inches(8), Inches(1.5), Pt(48), CYAN, bold=True)
add_text(slide1, "Integração entre pesquisa química e IA Generativa (Ollama/Gemma2)", Inches(0.8), Inches(4.7), Inches(8), Inches(0.5), Pt(16), WHITE)
add_text(slide1, "Nomes dos Integrantes | Instituição de Ensino", Inches(0.8), Inches(6.5), Inches(8), Inches(0.5), Pt(14), GRAY)

# ---------------------------------------------------------
# SLIDE 2: Introdução
# ---------------------------------------------------------
slide2 = prs.slides.add_slide(blank_layout)
set_background(slide2, BG_COLOR)
add_text(slide2, "INTRODUÇÃO", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(2.8), Inches(4), CARD_BG)
add_text(slide2, "01", Inches(0.7), Inches(2.2), Inches(2), Inches(1), Pt(48), CYAN, bold=True)
add_text(slide2, "IA na Pesquisa", Inches(0.7), Inches(3.2), Inches(2.4), Inches(0.5), Pt(20), WHITE, bold=True)
add_text(slide2, "Uso de Inteligência Artificial para acelerar pesquisas e descobertas.", Inches(0.7), Inches(4), Inches(2.4), Inches(2), Pt(14), GRAY)

add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(2), Inches(2.8), Inches(4), CARD_BG)
add_text(slide2, "02", Inches(3.8), Inches(2.2), Inches(2), Inches(1), Pt(48), CYAN, bold=True)
add_text(slide2, "Tradução Complexa", Inches(3.8), Inches(3.2), Inches(2.4), Inches(0.5), Pt(20), WHITE, bold=True)
add_text(slide2, "Tradução de nomes populares para estruturas químicas complexas (SMILES).", Inches(3.8), Inches(4), Inches(2.4), Inches(2), Pt(14), GRAY)

add_shape(slide2, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(2), Inches(2.8), Inches(4), CARD_BG)
add_text(slide2, "03", Inches(6.9), Inches(2.2), Inches(2), Inches(1), Pt(48), CYAN, bold=True)
add_text(slide2, "Validação Oficial", Inches(6.9), Inches(3.2), Inches(2.4), Inches(0.5), Pt(20), WHITE, bold=True)
add_text(slide2, "Integração direta com bases oficiais como NCBI e PubChem para validação.", Inches(6.9), Inches(4), Inches(2.4), Inches(2), Pt(14), GRAY)

# ---------------------------------------------------------
# SLIDE 3: O Problema
# ---------------------------------------------------------
slide3 = prs.slides.add_slide(blank_layout)
set_background(slide3, BG_COLOR)
add_text(slide3, "O PROBLEMA", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

add_shape(slide3, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(3.8), Inches(4), CARD_BG)
add_text(slide3, "Por que buscar compostos\né tão difícil?", Inches(0.8), Inches(3.2), Inches(3.4), Inches(2), Pt(28), WHITE, bold=True)

# Linha do tempo de dores (Pain points timeline)
add_shape(slide3, MSO_SHAPE.OVAL, Inches(4.8), Inches(2.4), Inches(0.2), Inches(0.2), CYAN)
add_shape(slide3, MSO_SHAPE.RECTANGLE, Inches(4.88), Inches(2.6), Inches(0.04), Inches(1.2), CYAN)
add_shape(slide3, MSO_SHAPE.OVAL, Inches(4.8), Inches(3.8), Inches(0.2), Inches(0.2), CYAN)
add_shape(slide3, MSO_SHAPE.RECTANGLE, Inches(4.88), Inches(4.0), Inches(0.04), Inches(1.2), CYAN)
add_shape(slide3, MSO_SHAPE.OVAL, Inches(4.8), Inches(5.2), Inches(0.2), Inches(0.2), CYAN)

add_text(slide3, "Bases tradicionais exigem nomenclatura IUPAC exata.", Inches(5.2), Inches(2.3), Inches(4.5), Inches(1), Pt(18), WHITE)
add_text(slide3, "Pesquisadores perdem tempo traduzindo nomes populares.", Inches(5.2), Inches(3.7), Inches(4.5), Inches(1), Pt(18), WHITE)
add_text(slide3, "Falta de ferramentas locais e centralizadas.", Inches(5.2), Inches(5.1), Inches(4.5), Inches(1), Pt(18), WHITE)

# ---------------------------------------------------------
# SLIDE 4: Objetivos
# ---------------------------------------------------------
slide4 = prs.slides.add_slide(blank_layout)
set_background(slide4, BG_COLOR)
add_text(slide4, "OBJETIVOS DO PROJETO", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(2), CYAN)
add_text(slide4, "Objetivo Geral", Inches(0.8), Inches(1.7), Inches(8), Inches(0.5), Pt(16), BG_COLOR, bold=True)
add_text(slide4, "Desenvolver API REST e Interface Web para consulta e comparação\nde moléculas orquestrando IA Generativa e PubChem.", Inches(0.8), Inches(2.2), Inches(8), Inches(1), Pt(22), BG_COLOR, bold=True)

add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(2.8), Inches(2.5), CARD_BG)
add_text(slide4, "Implementar IA Local", Inches(0.7), Inches(4.3), Inches(2.4), Inches(0.5), Pt(16), CYAN, bold=True)
add_text(slide4, "Utilizar Ollama com modelo Gemma2 para gerar SMILES com segurança e privacidade.", Inches(0.7), Inches(4.8), Inches(2.4), Inches(1.5), Pt(14), GRAY)

add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(4), Inches(2.8), Inches(2.5), CARD_BG)
add_text(slide4, "Garantir Segurança", Inches(3.8), Inches(4.3), Inches(2.4), Inches(0.5), Pt(16), CYAN, bold=True)
add_text(slide4, "Proteger acesso com autenticação JWT e usar banco de dados relacional isolado.", Inches(3.8), Inches(4.8), Inches(2.4), Inches(1.5), Pt(14), GRAY)

add_shape(slide4, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(4), Inches(2.8), Inches(2.5), CARD_BG)
add_text(slide4, "Performance e Métricas", Inches(6.9), Inches(4.3), Inches(2.4), Inches(0.5), Pt(16), CYAN, bold=True)
add_text(slide4, "Mensurar o tempo de resposta entre a IA (local) e a API pública (nuvem).", Inches(6.9), Inches(4.8), Inches(2.4), Inches(1.5), Pt(14), GRAY)

# ---------------------------------------------------------
# SLIDE 5: Público-Alvo
# ---------------------------------------------------------
slide5 = prs.slides.add_slide(blank_layout)
set_background(slide5, BG_COLOR)
add_text(slide5, "PÚBLICO-ALVO", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

targets = [
    ("Pesquisadores", "Iniciantes em\nIA e Quimioinformática."),
    ("Estudantes", "Química Farmacêutica\ne áreas correlatas."),
    ("Laboratórios", "Buscam prototipação\nrápida e barata.")
]
for i in range(3):
    add_shape(slide5, MSO_SHAPE.OVAL, Inches(1.2 + i*2.8), Inches(2.5), Inches(2), Inches(2), BG_COLOR, line_color=CYAN, hollow=True)
    add_text(slide5, targets[i][0], Inches(0.5 + i*2.8), Inches(3.2), Inches(3.4), Inches(0.5), Pt(18), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide5, targets[i][1], Inches(0.5 + i*2.8), Inches(5), Inches(3.4), Inches(1), Pt(16), GRAY, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 6: Tecnologias
# ---------------------------------------------------------
slide6 = prs.slides.add_slide(blank_layout)
set_background(slide6, BG_COLOR)
add_text(slide6, "TECNOLOGIAS UTILIZADAS", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

techs = [
    ("BACKEND", "Python, FastAPI\nSQLAlchemy", 0.5, 2),
    ("FRONTEND", "Streamlit\n(Python)", 3.6, 2),
    ("DATABASE", "SQLite 3", 6.7, 2),
    ("IA CORE", "Ollama\nModelo Gemma2", 0.5, 4.5),
    ("APIs EXTERNAS", "PubChem\nREST API", 3.6, 4.5),
    ("SEGURANÇA", "PyJWT\nbcrypt", 6.7, 4.5)
]

for t in techs:
    add_shape(slide6, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(t[2]), Inches(t[3]), Inches(2.8), Inches(2), CARD_BG)
    # Highlight tag
    add_shape(slide6, MSO_SHAPE.RECTANGLE, Inches(t[2]), Inches(t[3]+0.3), Inches(0.1), Inches(0.6), CYAN)
    add_text(slide6, t[0], Inches(t[2]+0.3), Inches(t[3]+0.3), Inches(2.4), Inches(0.5), Pt(14), CYAN, bold=True)
    add_text(slide6, t[1], Inches(t[2]+0.3), Inches(t[3]+0.7), Inches(2.4), Inches(1), Pt(22), WHITE, bold=True)

# ---------------------------------------------------------
# SLIDE 7: Arquitetura do Sistema
# ---------------------------------------------------------
slide7 = prs.slides.add_slide(blank_layout)
set_background(slide7, BG_COLOR)
add_text(slide7, "ARQUITETURA DO SISTEMA", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

# Blocos de diagrama
add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.2), Inches(2.5), CARD_BG)
add_text(slide7, "Streamlit\n(Interface)", Inches(0.5), Inches(3.3), Inches(2.2), Inches(1), Pt(18), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide7, MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(3.6), Inches(0.6), Inches(0.2), CYAN)

add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2), Inches(3), Inches(3.5), CARD_BG)
add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2), Inches(3), Inches(0.4), CYAN)
add_text(slide7, "FASTAPI (Orquestrador)", Inches(3.5), Inches(2.05), Inches(3), Inches(1), Pt(14), BG_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide7, "Delega orquestração\nassíncrona e\nunifica os dados.", Inches(3.5), Inches(3.2), Inches(3), Inches(1), Pt(16), WHITE, alignment=PP_ALIGN.CENTER)

add_shape(slide7, MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(2.5), Inches(0.6), Inches(0.15), CYAN)
add_shape(slide7, MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(3.6), Inches(0.6), Inches(0.15), CYAN)
add_shape(slide7, MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(4.7), Inches(0.6), Inches(0.15), CYAN)

add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(2.1), Inches(2.2), Inches(0.8), BG_COLOR, line_color=CYAN, hollow=True)
add_text(slide7, "Ollama / Gemma2", Inches(7.3), Inches(2.35), Inches(2.2), Inches(0.5), Pt(14), CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(3.2), Inches(2.2), Inches(0.8), BG_COLOR, line_color=CYAN, hollow=True)
add_text(slide7, "PubChem API", Inches(7.3), Inches(3.45), Inches(2.2), Inches(0.5), Pt(14), CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(4.3), Inches(2.2), Inches(0.8), BG_COLOR, line_color=CYAN, hollow=True)
add_text(slide7, "SQLite DB", Inches(7.3), Inches(4.55), Inches(2.2), Inches(0.5), Pt(14), CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 8: Demonstração Frontend
# ---------------------------------------------------------
slide8 = prs.slides.add_slide(blank_layout)
set_background(slide8, BG_COLOR)
add_text(slide8, "DEMONSTRAÇÃO: FRONTEND", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

# Mockup estilo Browser
add_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), CARD_BG)
# Barra do Browser
add_shape(slide8, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4), RGBColor(10, 15, 20))
add_shape(slide8, MSO_SHAPE.OVAL, Inches(0.7), Inches(1.65), Inches(0.15), Inches(0.15), RGBColor(255, 95, 86)) # Red
add_shape(slide8, MSO_SHAPE.OVAL, Inches(0.95), Inches(1.65), Inches(0.15), Inches(0.15), RGBColor(255, 189, 46)) # Yellow
add_shape(slide8, MSO_SHAPE.OVAL, Inches(1.2), Inches(1.65), Inches(0.15), Inches(0.15), RGBColor(39, 201, 63)) # Green
add_shape(slide8, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.6), Inches(6), Inches(0.2), CARD_BG)

# Área pontilhada para o print
add_shape(slide8, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(8), Inches(4.3), CARD_BG, line_color=GRAY, hollow=True)
add_text(slide8, "[ ARRASTE SEU PRINT DO STREAMLIT AQUI ]", Inches(1), Inches(4.2), Inches(8), Inches(1), Pt(20), GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 9: Demonstração Backend & DER
# ---------------------------------------------------------
slide9 = prs.slides.add_slide(blank_layout)
set_background(slide9, BG_COLOR)
add_text(slide9, "DEMONSTRAÇÃO: DER E BANCO DE DADOS", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

# Mockup estilo Ferramenta/IDE
add_shape(slide9, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5), CARD_BG)
add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(2.5), Inches(5.5), RGBColor(20, 24, 34))

add_text(slide9, "SCHEMA DB", Inches(0.7), Inches(1.8), Inches(2), Inches(0.5), Pt(14), CYAN, bold=True)
add_text(slide9, "• users (1)\n• histories (N)\n• molecules (1)\n• logs (N)\n• sessions (N)", Inches(0.7), Inches(2.3), Inches(2), Inches(2), Pt(14), WHITE)

add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(1.7), Inches(6.1), Inches(5.1), CARD_BG, line_color=CYAN, hollow=True)
add_text(slide9, "[ ARRASTE A IMAGEM DO DER.png AQUI ]", Inches(3.2), Inches(4.2), Inches(6.1), Inches(1), Pt(18), CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 10: Fluxo Completo
# ---------------------------------------------------------
slide10 = prs.slides.add_slide(blank_layout)
set_background(slide10, BG_COLOR)
add_text(slide10, "FLUXO DE INTEGRAÇÃO (PONTA A PONTA)", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

add_shape(slide10, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4), Inches(9), Inches(0.05), CYAN)

steps = [
    ("Login & Input", "Usuário pesquisa\n'Aspirina'"),
    ("Backend API", "Endpoint\n/molecules/search"),
    ("Orquestração", "LLM traduz e\nPubChem valida"),
    ("Persistência", "Banco salva a\nmétrica de tempo"),
    ("Renderização", "Frontend exibe\no visual 2D")
]

for i in range(5):
    # Circle node
    add_shape(slide10, MSO_SHAPE.OVAL, Inches(1 + i*1.75), Inches(3.85), Inches(0.35), Inches(0.35), BG_COLOR, line_color=CYAN)
    add_text(slide10, f"0{i+1}", Inches(0.8 + i*1.75), Inches(3.3), Inches(0.75), Inches(0.5), Pt(24), CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide10, steps[i][0], Inches(0.2 + i*1.75), Inches(4.4), Inches(2), Inches(0.5), Pt(14), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide10, steps[i][1], Inches(0.2 + i*1.75), Inches(4.8), Inches(2), Inches(1), Pt(12), GRAY, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------
# SLIDE 11: Resultados e Conclusão
# ---------------------------------------------------------
slide11 = prs.slides.add_slide(blank_layout)
set_background(slide11, BG_COLOR)
add_text(slide11, "RESULTADOS E FUTURO", Inches(0.5), Inches(0.5), Inches(9), Inches(1), Pt(32), CYAN, bold=True)

add_shape(slide11, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(2.2), CYAN)
add_text(slide11, "Missão Cumprida: Ecossistema Funcional", Inches(0.8), Inches(2), Inches(8.4), Inches(0.5), Pt(28), BG_COLOR, bold=True)
add_text(slide11, "A integração bem-sucedida atestou a viabilidade técnica e prática do uso de IAs locais (Gemma2) para agilizar e democratizar fluxos de quimioinformática.", Inches(0.8), Inches(2.7), Inches(8.4), Inches(1), Pt(18), BG_COLOR)

add_shape(slide11, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(4.3), Inches(2.4), CARD_BG)
add_text(slide11, "Impacto Imediato", Inches(0.8), Inches(4.5), Inches(3.7), Inches(0.5), Pt(20), CYAN, bold=True)
add_text(slide11, "• Redução do tempo de tradução química.\n• Centralização das bases de dados.\n• Interface amigável para iniciantes.", Inches(0.8), Inches(5), Inches(3.7), Inches(1.5), Pt(16), WHITE)

add_shape(slide11, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(4.3), Inches(4.3), Inches(2.4), CARD_BG)
add_text(slide11, "Próximos Passos", Inches(5.5), Inches(4.5), Inches(3.7), Inches(0.5), Pt(20), CYAN, bold=True)
add_text(slide11, "• Predição avançada de propriedades moleculares.\n• Suporte nativo a múltiplos LLMs.\n• Expandir integrações com outras APIs científicas.", Inches(5.5), Inches(5), Inches(3.7), Inches(1.5), Pt(16), WHITE)

prs.save('Apresentacao_Premium_Quimiochat.pptx')
print("Apresentacao_Premium_Quimiochat.pptx gerada com sucesso.")
