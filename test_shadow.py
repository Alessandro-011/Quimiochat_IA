from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2))
try:
    shadow = shape.shadow
    shadow.inherit = False
    print("Shadow attribute exists!")
except Exception as e:
    print(f"Error: {e}")
prs.save("test_shadow.pptx")
