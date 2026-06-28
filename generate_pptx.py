from pptx import Presentation
from pptx.util import Inches, Pt

# Criação da apresentação
prs = Presentation()

# Layouts padrão do python-pptx (0 = Título, 1 = Título e Conteúdo)
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# ---------------------------------------------------------
# SLIDE 1: Capa [cite: 204]
# ---------------------------------------------------------
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sistema de Quimioinformática Inteligente"
subtitle.text = "Integração entre pesquisa química e IA Generativa (Ollama/Gemma2)\n\nNome dos Integrantes\nInstituição de Ensino"

# ---------------------------------------------------------
# SLIDE 2: Introdução ao Tema [cite: 205, 239, 240, 241, 242]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "1. Introdução"
tf = body.text_frame
tf.text = "A intersecção entre Química e Tecnologia"
tf.add_paragraph().text = "Uso de Inteligência Artificial para acelerar pesquisas."
tf.add_paragraph().text = "Tradução de nomes populares para estruturas químicas complexas (SMILES)."
tf.add_paragraph().text = "Integração com bases oficiais (NCBI/PubChem) para validação."

# ---------------------------------------------------------
# SLIDE 3: O Problema [cite: 206, 243, 244, 245, 246, 247]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "2. O Problema"
tf = body.text_frame
tf.text = "Dificuldade na busca de compostos químicos"
tf.add_paragraph().text = "Bases tradicionais exigem nomenclatura IUPAC exata."
tf.add_paragraph().text = "Pesquisadores iniciantes perdem tempo traduzindo nomes populares (ex: Aspirina) para fórmulas."
tf.add_paragraph().text = "Falta de ferramentas locais rápidas que centralizem IA e dados oficiais."

# ---------------------------------------------------------
# SLIDE 4: Objetivos do Projeto [cite: 207, 248, 249, 250, 251, 252]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "3. Objetivos"
tf = body.text_frame
tf.text = "Objetivo Geral:"
tf.add_paragraph().text = "Desenvolver uma API REST e uma interface web para consulta e comparação de moléculas usando IA e PubChem."
p = tf.add_paragraph()
p.text = "Objetivos Específicos:"
tf.add_paragraph().text = "Implementar LLM local (Ollama/Gemma2) para gerar SMILES."
tf.add_paragraph().text = "Garantir segurança com JWT e banco de dados relacional."
tf.add_paragraph().text = "Mensurar tempo de resposta entre a IA e a API pública."

# ---------------------------------------------------------
# SLIDE 5: Público-Alvo [cite: 208]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "4. Público-Alvo"
tf = body.text_frame
tf.text = "Pesquisadores iniciantes na área de Inteligência Artificial e Quimioinformática."
tf.add_paragraph().text = "Estudantes de química farmacêutica e áreas correlatas."
tf.add_paragraph().text = "Laboratórios que necessitam de ferramentas rápidas de prototipação."

# ---------------------------------------------------------
# SLIDE 6: Tecnologias Utilizadas [cite: 209, 253, 254, 255, 256, 257, 258, 259, 260]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "5. Tecnologias Utilizadas"
tf = body.text_frame
tf.text = "Backend: Python, FastAPI, SQLAlchemy, Pydantic"
tf.add_paragraph().text = "Frontend: Streamlit"
tf.add_paragraph().text = "Banco de Dados: SQLite 3"
tf.add_paragraph().text = "Inteligência Artificial: Ollama (Modelo Gemma2)"
tf.add_paragraph().text = "APIs Externas: PubChem REST API"
tf.add_paragraph().text = "Segurança: PyJWT, bcrypt"

# ---------------------------------------------------------
# SLIDE 7: Estrutura do Sistema (Arquitetura) [cite: 210, 261, 262, 263, 264, 265, 266]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "6. Arquitetura do Sistema"
tf = body.text_frame
tf.text = "Padrão adotado: MVC e microsserviços lógicos."
tf.add_paragraph().text = "Interface (Streamlit) envia JSON via HTTP para o Backend."
tf.add_paragraph().text = "Rotas do FastAPI delegam a orquestração assíncrona."
tf.add_paragraph().text = "Consultas paralelas disparadas para Ollama e PubChem."
tf.add_paragraph().text = "Dados unificados e salvos no SQLite antes de retornar ao usuário."

# ---------------------------------------------------------
# SLIDE 8: Demonstração - Frontend [cite: 211, 267, 268, 269, 270, 271, 272, 273]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "7. Demonstração: Interface Frontend"
tf = body.text_frame
tf.text = "[ESPAÇO RESERVADO PARA PRINT DO STREAMLIT]"
tf.add_paragraph().text = "-> Cole aqui um print da tela de Login."
tf.add_paragraph().text = "-> Cole aqui um print da tela de Pesquisa mostrando as moléculas."

# ---------------------------------------------------------
# SLIDE 9: Demonstração - Backend e DB [cite: 212, 274, 275, 276, 277, 278, 279, 280, 281]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "8. Demonstração: Backend e Banco de Dados"
tf = body.text_frame
tf.text = "[ESPAÇO RESERVADO PARA O DER]"
tf.add_paragraph().text = "-> Cole aqui a imagem DER.png gerada na Atividade 2."
tf.add_paragraph().text = "O banco conta com 5 tabelas interligadas gerenciando histórico de pesquisas (1:N) e resultados (1:1)."

# ---------------------------------------------------------
# SLIDE 10: Integração Completa [cite: 213, 282, 283, 284, 285, 286, 287, 288]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "9. Integração (O Fluxo de Ponta a Ponta)"
tf = body.text_frame
tf.text = "1. Usuário autentica e pesquisa 'Aspirina'."
tf.add_paragraph().text = "2. Requisição bate no endpoint /molecules/search."
tf.add_paragraph().text = "3. LLM traduz e PubChem valida simultaneamente."
tf.add_paragraph().text = "4. Banco de dados salva a métrica de tempo."
tf.add_paragraph().text = "5. Frontend renderiza o visual 2D comparativo."

# ---------------------------------------------------------
# SLIDE 11: Resultados e Conclusão [cite: 214, 215, 289, 290, 291, 292, 293, 294]
# ---------------------------------------------------------
slide = prs.slides.add_slide(bullet_slide_layout)
title, body = slide.shapes.title, slide.placeholders[1]
title.text = "10. Resultados e Considerações Finais"
tf = body.text_frame
tf.text = "Entregamos um ecossistema funcional completo."
tf.add_paragraph().text = "Integração bem-sucedida atestou viabilidade de IAs locais para fluxos químicos."
tf.add_paragraph().text = "Melhorias Futuras: Implementar predição de propriedades moleculares e suporte a outros LLMs."

# Salvar o arquivo
prs.save('Apresentacao_TCC_Quimiochat.pptx')
print("✅ Sucesso! O arquivo Apresentacao_TCC_Quimiochat.pptx foi gerado na pasta.")
